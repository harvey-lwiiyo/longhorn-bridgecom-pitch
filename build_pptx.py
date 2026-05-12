"""
Longhorn Insurance × BridgeCom Narratives — Pitch Deck PPTX generator
Produces a fully editable PowerPoint from the design content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x00, 0x6E)
RED    = RGBColor(0xC8, 0x12, 0x1E)
CREAM  = RGBColor(0xF6, 0xF4, 0xEF)
CRMWRM = RGBColor(0xEC, 0xE7, 0xDC)
INK    = RGBColor(0x15, 0x15, 0x2A)
MUTED  = RGBColor(0x5B, 0x5B, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xBA, 0x75, 0x17)
GREEN  = RGBColor(0x27, 0x50, 0x0A)
RULE   = RGBColor(0xD8, 0xD7, 0xE0)

# ── Fonts ─────────────────────────────────────────────────────────────────────
SERIF = "Georgia"        # Instrument Serif substitute
SANS  = "Calibri"        # Geist substitute
MONO  = "Courier New"    # Geist Mono substitute

# ── Layout constants (inches) ─────────────────────────────────────────────────
PX = 13.333 / 1920      # 1px in inches
PAD_X  = 110 * PX       # 0.764"
PAD_T  = 100 * PX       # 0.694"
PAD_B  = 90  * PX       # 0.625"
CW     = 13.333 - 2 * PAD_X   # content width
FOOTER_Y = SH.inches - PAD_B - 0.22

# ── Helpers ───────────────────────────────────────────────────────────────────

def px(n): return Inches(n * PX)

def bg(slide, color):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0.75)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        font=SANS, size=18, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True, fill=None, line_color=None,
        valign=None, spacing_before=0, spacing_after=0):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_before:
        p.space_before = Pt(spacing_before)
    if spacing_after:
        p.space_after = Pt(spacing_after)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    if fill:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if line_color:
        tb.line.color.rgb = line_color
    return tb

def multiline_txt(slide, lines, x, y, w, h,
                  font=SANS, size=18, color=INK, align=PP_ALIGN.LEFT,
                  line_spacing=1.15):
    """Add a text box with multiple paragraphs (list of (text, bold, italic, color, size))."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, False, False, color, size)
        text, bold, italic, col, sz = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(sz)
        run.font.color.rgb = col if col else color
        run.font.bold = bold
        run.font.italic = italic
    return tb

def hline(slide, x, y, w, color=RULE, width=Pt(0.5)):
    """Horizontal rule."""
    conn = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    conn.line.color.rgb = color
    conn.line.width = width

def corner_lockup(slide, section_label, dark=False):
    """BridgeCom logo lockup (top-left) and section label (top-right)."""
    c = WHITE if dark else MUTED
    ci = WHITE if dark else INK
    # Brand name
    txt(slide, "BridgeCom Narratives", PAD_X, 0.35, 3.5, 0.35,
        font=SANS, size=10, color=ci, bold=True)
    txt(slide, "FOR LONGHORN INSURANCE", PAD_X, 0.65, 3.5, 0.22,
        font=MONO, size=7.5, color=c)
    # Section label right-aligned
    txt(slide, section_label,
        13.333 - PAD_X - 2.5, 0.35, 2.5, 0.35,
        font=MONO, size=9, color=c, align=PP_ALIGN.RIGHT)

def footer(slide, left="Longhorn × BridgeCom", right="", dark=False):
    """Footer bar."""
    c = RGBColor(0x90, 0x90, 0xA8) if dark else MUTED
    hline(slide, PAD_X, FOOTER_Y + 0.14, CW, color=RGBColor(0x80,0x80,0x90) if dark else RULE)
    txt(slide, left, PAD_X, FOOTER_Y, 4, 0.25, font=MONO, size=8, color=c)
    if right:
        txt(slide, right, 13.333 - PAD_X - 2, FOOTER_Y, 2, 0.25,
            font=MONO, size=8, color=c, align=PP_ALIGN.RIGHT)

def eyebrow(slide, text, x, y, color=MUTED, dark=False):
    """Eyebrow label with red dot."""
    col = RGBColor(0xA0,0xA0,0xC0) if dark else MUTED
    txt(slide, "• " + text.upper(), x, y, CW, 0.28,
        font=MONO, size=9, color=col)

def divider_slide(prs, part_num, part_name, subtitle):
    """Full navy section-divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    # Giant ordinal watermark
    txt(slide, part_num,
        PAD_X - 0.3, 1.2, 4, 3.5,
        font=SERIF, size=200, color=RGBColor(0x28,0x20,0x9A),
        italic=True)
    # Part label
    eyebrow(slide, f"Part {part_num}", PAD_X + 3.2, 2.5, dark=True)
    # Title
    txt(slide, f"The {part_name}", PAD_X + 3.2, 2.85, 8, 1.5,
        font=SERIF, size=80, color=WHITE, italic=True)
    # Subtitle
    txt(slide, subtitle, PAD_X + 3.2, 4.0, 7.5, 0.9,
        font=SANS, size=20, color=RGBColor(0xC0,0xC0,0xE0))
    corner_lockup(slide, f"Part {part_num}", dark=True)
    footer(slide, "BridgeCom Narratives · 2026", "", dark=True)
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]  # blank layout


# ─────────────────────────────────────────────────────────────────────────────
# 01 · COVER
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
# Red gradient accent block (right side)
rect(s, 9.5, 3.5, 4, 5, fill=RGBColor(0x80, 0x05, 0x10))
rect(s, 10.5, 0, 3, 8, fill=RGBColor(0x06, 0x00, 0x55))
# Cover lockup
txt(s, "BRIDGECOM NARRATIVES", PAD_X, 1.0, 6, 0.35,
    font=MONO, size=10, color=RGBColor(0xCC,0xCC,0xFF))
txt(s, "×", PAD_X + 6.1, 0.95, 0.5, 0.4,
    font=SERIF, size=22, color=RGBColor(0x80,0x80,0xCC), italic=True)
txt(s, "LONGHORN INSURANCE BROKERS", PAD_X + 6.6, 1.0, 5.5, 0.35,
    font=MONO, size=10, color=RGBColor(0xCC,0xCC,0xFF))
# Main display title
txt(s, "Brand Growth", PAD_X, 1.55, 9, 1.5,
    font=SERIF, size=108, color=WHITE, italic=False)
txt(s, "Pitch", PAD_X, 2.9, 9, 1.4,
    font=SERIF, size=108, color=RED, italic=True)
# Dividing line
hline(s, PAD_X, 4.45, 7, color=RGBColor(0x40,0x40,0x80), width=Pt(1))
# Meta
txt(s, "A social media strategy to build Longhorn's presence,\ntrust, and pipeline from zero — across five channels.",
    PAD_X, 4.6, 8, 0.8, font=SANS, size=16, color=RGBColor(0xB0,0xB0,0xD8))
# Bottom meta row
txt(s, "BridgeCom Narratives  ·  Lusaka  ·  May 2026",
    PAD_X, 6.75, 7, 0.3, font=MONO, size=9, color=RGBColor(0x80,0x80,0xA8))


# ─────────────────────────────────────────────────────────────────────────────
# 02 · SECTION: THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
divider_slide(prs, "01", "Problem",
              "Five platforms. One number on each. The market does not know Longhorn exists.")


# ─────────────────────────────────────────────────────────────────────────────
# 03 · THE TRUTH
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Problem · 01")
eyebrow(s, "The truth, in one sentence", PAD_X, PAD_T)
txt(s, "You do good work.\nThe market does not know you exist.",
    PAD_X, PAD_T + 0.38, CW, 2.2,
    font=SERIF, size=54, color=INK, italic=True)
hline(s, PAD_X, PAD_T + 2.75, 0.7, color=INK, width=Pt(2))
txt(s, "Longhorn has product, licence, expertise, and a real broker model that legally serves the client, "
       "not the institution. None of that is visible online. A buyer searching for an insurance broker in "
       "Lusaka today will not find Longhorn on Google, on LinkedIn, on Facebook, on Instagram, or on WhatsApp. "
       "They will find your competitors.",
    PAD_X, PAD_T + 3.05, 9.5, 1.5,
    font=SANS, size=16, color=MUTED)
footer(s, "Longhorn × BridgeCom", "03 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 04 · FOOTPRINT TODAY
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Problem · 02")
eyebrow(s, "Your digital footprint, May 2026", PAD_X, PAD_T)
txt(s, "Five platforms.  One number on each.",
    PAD_X, PAD_T + 0.38, CW, 0.8,
    font=SERIF, size=42, color=INK, italic=False)

platforms = [
    ("LinkedIn",  "Company page followers · Zero leadership posts · Zero Lead Gen forms"),
    ("Facebook",  "Page followers · No live page · No retail or SME audience built"),
    ("Instagram", "Followers · No Reels · No 22–38 professional audience reached"),
    ("Google",    "Ranking keywords · No landing pages · Invisible to anyone searching"),
    ("WhatsApp",  "Broadcast subscribers · No automated reply · No nurture sequence"),
]
card_w = CW / 5 - 0.08
for i, (plat, metric) in enumerate(platforms):
    cx = PAD_X + i * (card_w + 0.10)
    cy = PAD_T + 1.35
    rect(s, cx, cy, card_w, 3.4, fill=WHITE, line=RULE, line_w=Pt(0.5))
    # Dashed border simulation (solid thin border)
    txt(s, plat.upper(), cx + 0.18, cy + 0.18, card_w - 0.3, 0.28,
        font=MONO, size=8, color=INK, bold=True)
    txt(s, "0", cx + 0.18, cy + 0.45, card_w - 0.3, 1.1,
        font=SERIF, size=68, color=RED)
    txt(s, metric, cx + 0.18, cy + 1.55, card_w - 0.3, 1.4,
        font=SANS, size=9.5, color=MUTED, wrap=True)

txt(s, "This is not a small problem to fix later. Every channel above is where a decision-maker "
       "in Lusaka decides who to call, who to trust, and who to pay. Longhorn is absent from all of them.",
    PAD_X, 5.45, CW, 0.55, font=SANS, size=13, color=MUTED)
footer(s, "Longhorn × BridgeCom", "04 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 05 · WHERE ZAMBIA DECIDES
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Problem · 03")
eyebrow(s, "Where Zambia actually decides", PAD_X, PAD_T)
# Left
txt(s, "The decision happens\nbefore they call you.",
    PAD_X, PAD_T + 0.38, 4.2, 1.3,
    font=SERIF, size=36, color=INK, italic=True)
txt(s, "98% of Zambians use WhatsApp daily. Facebook reaches 2.4M users. "
       "LinkedIn has 180k+ professionals. Google processes 300M+ searches monthly in Zambia. "
       "Every one of these platforms is a place where someone decides if Longhorn is credible "
       "before they ever pick up the phone.",
    PAD_X, PAD_T + 1.9, 4.2, 2.2,
    font=SANS, size=13, color=MUTED)

# Right — channel bars
bars = [
    ("WhatsApp",         98, True),
    ("Facebook",         62, False),
    ("Google Search",    55, False),
    ("Instagram",        38, True),
    ("LinkedIn",         18, False),
]
bar_x = PAD_X + 4.6
bar_y = PAD_T + 0.3
bar_w_total = CW - 4.6
for i, (name, pct, highlight) in enumerate(bars):
    row_y = bar_y + i * 1.02
    txt(s, name, bar_x, row_y, 1.7, 0.32, font=SANS, size=14, color=INK, bold=True)
    # Track background
    rect(s, bar_x + 1.8, row_y + 0.02, bar_w_total - 2.6, 0.3,
         fill=RGBColor(0xE0,0xDE,0xD8))
    # Fill
    fill_w = (pct / 100) * (bar_w_total - 2.6)
    fc = RED if highlight else NAVY
    rect(s, bar_x + 1.8, row_y + 0.02, fill_w, 0.3, fill=fc)
    # Percentage
    txt(s, f"{pct}%", bar_x + bar_w_total - 0.7, row_y, 0.65, 0.32,
        font=SERIF, size=22, color=INK, align=PP_ALIGN.RIGHT)

footer(s, "Longhorn × BridgeCom", "05 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 06 · CATEGORY TRAP
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CRMWRM)
corner_lockup(s, "The Problem · 04")
eyebrow(s, "The trust problem", PAD_X, PAD_T)
txt(s, "Insurance is the most trusted category\nno one wants to think about.",
    PAD_X, PAD_T + 0.4, CW, 1.6,
    font=SERIF, size=52, color=INK, italic=True)

col1_x = PAD_X
col2_x = PAD_X + CW / 2 + 0.2
cy = PAD_T + 2.4

txt(s, "The category problem",
    col1_x, cy, CW / 2 - 0.3, 0.32,
    font=MONO, size=9, color=RED, bold=True)
txt(s, "People buy insurance when forced to, delay when not, and choose on price because "
       "they have no other signal. Without a visible brand, Longhorn competes on price by "
       "default — and will lose to anyone willing to cut margins.",
    col1_x, cy + 0.38, CW / 2 - 0.3, 1.5,
    font=SANS, size=14, color=MUTED)

txt(s, "The opportunity",
    col2_x, cy, CW / 2 - 0.3, 0.32,
    font=MONO, size=9, color=NAVY, bold=True)
txt(s, "A broker who is visible, human, and consistently useful before the sale is the only "
       "broker who gets to charge on value rather than price. That broker does not exist in "
       "Zambia yet. BridgeCom's job is to make Longhorn that broker.",
    col2_x, cy + 0.38, CW / 2 - 0.3, 1.5,
    font=SANS, size=14, color=MUTED)

footer(s, "Longhorn × BridgeCom", "06 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 07 · SECTION: THE SOLUTION
# ─────────────────────────────────────────────────────────────────────────────
divider_slide(prs, "02", "Solution",
              "A visible brand built on one structural truth: the broker works for you, not the insurer.")


# ─────────────────────────────────────────────────────────────────────────────
# 08 · THE REFRAME
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Solution · 01")
eyebrow(s, "The strategic reframe", PAD_X, PAD_T)
txt(s, "Stop selling insurance.\nStart selling certainty.",
    PAD_X, PAD_T + 0.38, CW, 1.0,
    font=SERIF, size=40, color=INK, italic=True)

# FROM box
rect(s, PAD_X, PAD_T + 1.65, CW / 2 - 0.3, 3.0, fill=WHITE, line=RULE, line_w=Pt(1))
txt(s, "FROM", PAD_X + 0.3, PAD_T + 1.85, 3, 0.25, font=MONO, size=9, color=MUTED)
reframe_from = [
    ("Insurance products", False),
    ("Policy features", False),
    ("Price comparisons", False),
    ("Claims ratios", False),
]
for i, (line, _) in enumerate(reframe_from):
    txt(s, f"– {line}", PAD_X + 0.3, PAD_T + 2.2 + i * 0.52, CW / 2 - 0.6, 0.4,
        font=SERIF, size=26, color=MUTED, italic=True)

# TO box
rect(s, PAD_X + CW / 2 + 0.1, PAD_T + 1.65, CW / 2 - 0.1, 3.0, fill=NAVY)
txt(s, "TO", PAD_X + CW / 2 + 0.4, PAD_T + 1.85, 3, 0.25, font=MONO, size=9,
    color=RGBColor(0xA0,0xA0,0xD0))
reframe_to = [
    ("Peace of mind, guaranteed",),
    ("What happens when it matters",),
    ("The broker who works for you",),
    ("Someone on your side. Always.",),
]
for i, (line,) in enumerate(reframe_to):
    txt(s, line, PAD_X + CW / 2 + 0.4, PAD_T + 2.2 + i * 0.52, CW / 2 - 0.6, 0.4,
        font=SERIF, size=26, color=WHITE, italic=True)

footer(s, "Longhorn × BridgeCom", "08 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 09 · ON YOUR SIDE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
# Shield graphic (simplified rectangle)
rect(s, 10.2, 1.0, 2.4, 5.0, fill=RGBColor(0x18,0x08,0x88))
rect(s, 10.4, 1.3, 2.0, 4.0, fill=RGBColor(0x10,0x04,0x70))
txt(s, "✓", 10.8, 2.8, 1.4, 1.4, font=SANS, size=72, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "PIA LICENSED", 10.3, 5.2, 2.5, 0.32,
    font=MONO, size=8, color=RGBColor(0x80,0x80,0xC0), align=PP_ALIGN.CENTER)

corner_lockup(s, "The Solution · 02", dark=True)
eyebrow(s, "The brand promise", PAD_X, PAD_T + 0.3, dark=True)
txt(s, "On your side.", PAD_X, PAD_T + 0.72, 9.5, 1.3,
    font=SERIF, size=110, color=WHITE)
txt(s, "Always.", PAD_X, PAD_T + 1.9, 9.5, 1.3,
    font=SERIF, size=110, color=RGBColor(0xE8,0xE6,0xF0), italic=True)
txt(s, "Every insurer works for their shareholders. Every bank works for its balance sheet. "
       "A broker is legally required to work for the person sitting across the table. "
       "That is not a marketing claim. It is a structural fact that no competitor can honestly match.",
    PAD_X, PAD_T + 3.35, 9.2, 1.2,
    font=SANS, size=15, color=RGBColor(0xC0,0xC0,0xE0))
footer(s, "Longhorn × BridgeCom", "09 / 24", dark=True)


# ─────────────────────────────────────────────────────────────────────────────
# 10 · THREE AUDIENCES
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Solution · 03")
eyebrow(s, "Three audiences, one brand", PAD_X, PAD_T)
txt(s, "Different people. Different fears.\nOne voice that works for all of them.",
    PAD_X, PAD_T + 0.38, CW, 0.9,
    font=SERIF, size=34, color=INK, italic=True)

audiences = [
    ("Corporate & SME",   NAVY, "HR Managers, CFOs, Business Owners.",
     "If my business burns down tomorrow, can I make payroll next month?"),
    ("Retail & Family",   RED,  "Working professionals, 25–50, Lusaka.",
     "If I die tomorrow, does my family keep the house?"),
    ("High Net Worth",    GREEN,"Business owners, senior professionals.",
     "If I lose one asset, does everything unravel?"),
]
card_w = CW / 3 - 0.1
for i, (title, color, who, fear) in enumerate(audiences):
    cx = PAD_X + i * (card_w + 0.15)
    cy = PAD_T + 1.5
    rect(s, cx, cy, card_w, 4.1, fill=WHITE, line=RULE, line_w=Pt(0.5))
    rect(s, cx, cy, card_w, 0.07, fill=color)  # top accent bar
    txt(s, title, cx + 0.22, cy + 0.2, card_w - 0.4, 0.5,
        font=SERIF, size=26, color=INK)
    txt(s, who, cx + 0.22, cy + 0.82, card_w - 0.4, 0.55,
        font=SANS, size=12, color=MUTED)
    hline(s, cx + 0.22, cy + 1.6, card_w - 0.44, color=RULE)
    txt(s, "CORE FEAR", cx + 0.22, cy + 1.72, card_w - 0.4, 0.22,
        font=MONO, size=7.5, color=RED)
    txt(s, f'"{fear}"', cx + 0.22, cy + 1.98, card_w - 0.4, 1.7,
        font=SERIF, size=16, color=INK, italic=True)

footer(s, "Longhorn × BridgeCom", "10 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 11 · PRODUCT REFRAMING
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Solution · 04")
eyebrow(s, "Product language, reframed", PAD_X, PAD_T)
txt(s, "Every product has a human name.",
    PAD_X, PAD_T + 0.38, CW, 0.6,
    font=SERIF, size=44, color=INK, italic=True)

# Table
headers = ["PRODUCT", "CATEGORY LANGUAGE", "LONGHORN LANGUAGE"]
col_widths = [2.2, 3.6, 5.7]
tx = PAD_X
ty = PAD_T + 1.2
# Header row
hline(s, PAD_X, ty, CW, color=INK, width=Pt(1.2))
for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
    txt(s, hdr, tx + sum(col_widths[:j]) + j*0.05, ty + 0.08, cw, 0.28,
        font=MONO, size=8.5, color=MUTED)
hline(s, PAD_X, ty + 0.38, CW, color=INK, width=Pt(1))

rows = [
    ("Motor",          "Comprehensive vehicle insurance",      "If someone writes it off — we sort it."),
    ("Life",           "Term life / endowment",                "Your family keeps the life you built."),
    ("Business",       "Commercial property & liability",      "Your business survives the unthinkable."),
    ("Medical",        "Medical expense cover",                "You get treated. Not declined."),
    ("Home",           "Household contents & structure",       "Everything inside still belongs to you."),
]
for r, (prod, old, new) in enumerate(rows):
    ry = ty + 0.38 + r * 0.72
    txt(s, prod, PAD_X + 0.05, ry + 0.1, col_widths[0] - 0.1, 0.5,
        font=SERIF, size=20, color=INK)
    txt(s, old, PAD_X + col_widths[0] + 0.15, ry + 0.1, col_widths[1] - 0.15, 0.5,
        font=SANS, size=14, color=MUTED, italic=True)
    txt(s, new, PAD_X + col_widths[0] + col_widths[1] + 0.25, ry + 0.1, col_widths[2] - 0.2, 0.5,
        font=SERIF, size=17, color=RED, italic=True)
    hline(s, PAD_X, ry + 0.64, CW, color=RULE, width=Pt(0.4))

footer(s, "Longhorn × BridgeCom", "11 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 12 · SECTION: THE PLAN
# ─────────────────────────────────────────────────────────────────────────────
divider_slide(prs, "03", "Plan",
              "Five channels. One funnel. Every touchpoint engineered to reduce resistance and increase action.")


# ─────────────────────────────────────────────────────────────────────────────
# 13 · ONE FUNNEL
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Plan · 01")
eyebrow(s, "Five channels, one funnel", PAD_X, PAD_T)
txt(s, "Every channel feeds the next.\nNothing is isolated.",
    PAD_X, PAD_T + 0.38, CW, 0.8,
    font=SERIF, size=40, color=INK, italic=True)

stages = [
    (NAVY,                      "01", "AWARENESS",   "LinkedIn + Facebook + Instagram",   "Reach the right person before they are looking"),
    (RGBColor(0x1A,0x10,0x8A),  "02", "INTEREST",    "Content + SEO + Organic",           "Give them a reason to keep paying attention"),
    (RGBColor(0x2A,0x18,0x80),  "03", "INTENT",      "Google Search + Landing Pages",     "Catch them when they are actively searching"),
    (RED,                        "04", "CONVERSION",  "WhatsApp + Broker contact",         "Move the conversation into a quote"),
    (RGBColor(0x8B,0x0A,0x16),  "05", "RETENTION",   "CRM + Broadcast + Referrals",       "Turn one client into many"),
]
fh = 0.72
for i, (color, num, stage, channels, desc) in enumerate(stages):
    fy = PAD_T + 1.42 + i * fh
    indent = i * 0.18
    rect(s, PAD_X + indent, fy, CW - indent, fh - 0.04, fill=color)
    txt(s, num, PAD_X + indent + 0.15, fy + 0.18, 0.5, 0.38,
        font=SERIF, size=28, color=RGBColor(0xFF,0xFF,0xFF,), italic=True)
    txt(s, stage, PAD_X + indent + 0.75, fy + 0.12, 2.0, 0.28,
        font=SANS, size=14, color=WHITE, bold=True)
    txt(s, channels, PAD_X + indent + 0.75, fy + 0.4, 2.5, 0.25,
        font=MONO, size=8, color=RGBColor(0xC0,0xC0,0xFF))
    txt(s, desc, PAD_X + indent + 3.5, fy + 0.2, CW - indent - 3.7, 0.36,
        font=SANS, size=12, color=RGBColor(0xE0,0xE0,0xFF))

footer(s, "Longhorn × BridgeCom", "13 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 14 · CHANNELS AT A GLANCE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Plan · 02")
eyebrow(s, "Each channel does one job", PAD_X, PAD_T)
txt(s, "No channel does everything. Each does one thing for the funnel.",
    PAD_X, PAD_T + 0.38, CW, 0.6,
    font=SERIF, size=36, color=INK, italic=True)

channels = [
    (NAVY, "LinkedIn",           "B2B · Decision-makers",
     "Gets in front of the person who signs the contract.",
     "Targets HR Managers, CFOs and Founders by job title, industry and company size in Lusaka.",
     "4–8", "leads / mo by M6"),
    (RED,  "Facebook & Instagram", "Retail · SME · Trust",
     "Builds the audience that warms before it buys.",
     "Reaches working Lusaka aged 25–50. Long-form trust on Facebook. Visual credibility on Instagram.",
     "10–20", "leads / mo by M6"),
    (GREEN,"Google",             "Intent · Capture",
     "Catches the person who is already looking.",
     "Three campaigns: broker intent, product intent, problem-aware. Lowest volume, highest conversion.",
     "3–6", "enquiries / mo by M6"),
    (GOLD, "WhatsApp",           "Conversion · The close",
     "Where the quote becomes a policy.",
     "Every other channel exists to move someone here. 98% open rate. 22-day nurture sequence.",
     "300+", "subscribers by M6"),
]
cw4 = CW / 4 - 0.1
for i, (color, name, badge, role, desc, num, lbl) in enumerate(channels):
    cx = PAD_X + i * (cw4 + 0.13)
    cy = PAD_T + 1.18
    rect(s, cx, cy, cw4, 4.55, fill=WHITE, line=RULE, line_w=Pt(0.5))
    rect(s, cx, cy, cw4, 0.06, fill=color)  # accent bar
    txt(s, badge.upper(), cx + 0.2, cy + 0.18, cw4 - 0.35, 0.25,
        font=MONO, size=7.5, color=MUTED)
    txt(s, name, cx + 0.2, cy + 0.45, cw4 - 0.35, 0.65,
        font=SERIF, size=28, color=INK, italic=True)
    txt(s, role, cx + 0.2, cy + 1.15, cw4 - 0.35, 0.55,
        font=SANS, size=11.5, color=INK, bold=True)
    txt(s, desc, cx + 0.2, cy + 1.72, cw4 - 0.35, 1.55,
        font=SANS, size=10.5, color=MUTED)
    txt(s, num, cx + 0.2, cy + 3.42, cw4 - 0.35, 0.62,
        font=SERIF, size=42, color=color)
    txt(s, lbl.upper(), cx + 0.2, cy + 4.05, cw4 - 0.35, 0.28,
        font=MONO, size=7, color=MUTED)

footer(s, "Longhorn × BridgeCom", "14 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 15 · BEHAVIOURAL DESIGN
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CRMWRM)
corner_lockup(s, "The Plan · 03")
eyebrow(s, "Behavioural design, applied at every interaction", PAD_X, PAD_T)
txt(s, "Every touchpoint reduces resistance.\nEvery touchpoint increases action.",
    PAD_X, PAD_T + 0.38, CW, 0.9,
    font=SERIF, size=36, color=INK, italic=True)

principles = [
    ("01", "Loss Aversion",
     "People are twice as motivated by avoiding a loss as gaining a benefit. "
     "Every headline leads with what is at stake — specific, local, vivid — before what is gained."),
    ("02", "Friction Reduction",
     "Four-field forms maximum. WhatsApp CTA on every page. LinkedIn Lead Gen pre-filled from profile. "
     "Every step reviewed: is anything here making the next move harder?"),
    ("03", "Social Proof",
     "Every claims proof story has a name, a business type, a location, and a specific outcome. "
     "'A Lusaka logistics company' beats 'a business in Zambia' every time."),
    ("04", "Trust Heuristics",
     "Named brokers, PIA licensing visible, consistent identity, SLA-fast response. "
     "People evaluate trust proxies because they cannot evaluate insurance quality before buying."),
    ("05", "Optimism Bias Fix",
     "SME owners assume the bad thing happens to other people. We make the worst case specific and real, "
     "then make the next step ridiculously easy to take."),
]
pw = CW / 5 - 0.1
for i, (n, h, d) in enumerate(principles):
    px2 = PAD_X + i * (pw + 0.12)
    py = PAD_T + 1.52
    rect(s, px2, py, pw, 4.22, fill=WHITE, line=RULE, line_w=Pt(0.5))
    txt(s, n, px2 + 0.18, py + 0.18, pw - 0.3, 0.65,
        font=SERIF, size=42, color=RED)
    txt(s, h, px2 + 0.18, py + 0.88, pw - 0.3, 0.55,
        font=SERIF, size=18, color=INK, bold=False)
    txt(s, d, px2 + 0.18, py + 1.5, pw - 0.3, 2.5,
        font=SANS, size=10, color=MUTED)

footer(s, "Longhorn × BridgeCom", "15 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 16 · SECTION: THE JOURNEY
# ─────────────────────────────────────────────────────────────────────────────
divider_slide(prs, "04", "Journey",
              "From invisible to compounding. Eighteen months. Four stages. One direction.")


# ─────────────────────────────────────────────────────────────────────────────
# 17 · 18-MONTH ROADMAP
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Journey · 01")
eyebrow(s, "Eighteen months, end to end", PAD_X, PAD_T)
txt(s, "From invisible to compounding. One stage at a time.",
    PAD_X, PAD_T + 0.38, CW, 0.6,
    font=SERIF, size=38, color=INK, italic=True)

stages_r = [
    (NAVY, "Stage 01", "Months 1–3",  "Foundation",
     "Build the infrastructure. Website, CRM, WhatsApp automation, social profiles, brand voice, analytics. "
     "No paid spend yet — there is nowhere for it to land.",
     "A credible brand online and a system to capture every enquiry."),
    (RED,  "Stage 02", "Months 4–7",  "Activation",
     "First paid campaigns. First Cover Check roadshow. First corporate risk briefings. "
     "SME workshop series. First two referral partners. Revenue begins here.",
     "10–25 enquiries per month. First policies signed from digital."),
    (GOLD, "Stage 03", "Months 8–12", "Optimisation",
     "Data drives every decision. Pause what fails, scale what works. CPL falls month over month. "
     "Pipeline reviewed weekly by leadership. SEO begins to land.",
     "30–50 enquiries per month. First full year-one KPI review."),
    (GREEN,"Stage 04", "Months 13–18","Compounding",
     "The system reinforces itself. Organic reach grows without proportional spend. "
     "Referrals multiply. Brand awareness reduces friction at every stage of the funnel.",
     "50+ enquiries per month. CPL down 30%+ vs Stage 2."),
]
sw4 = CW / 4 - 0.1
for i, (color, num, when, name, what, output) in enumerate(stages_r):
    cx = PAD_X + i * (sw4 + 0.13)
    cy = PAD_T + 1.2
    rect(s, cx, cy, sw4, 4.7, fill=WHITE, line=RULE, line_w=Pt(0.5))
    rect(s, cx, cy, sw4, 0.07, fill=color)
    txt(s, num.upper(), cx + 0.2, cy + 0.18, sw4 - 0.35, 0.22,
        font=MONO, size=8, color=MUTED)
    txt(s, when.upper(), cx + 0.2, cy + 0.42, sw4 - 0.35, 0.22,
        font=MONO, size=8, color=MUTED)
    txt(s, name, cx + 0.2, cy + 0.72, sw4 - 0.35, 0.52,
        font=SERIF, size=30, color=INK)
    txt(s, what, cx + 0.2, cy + 1.3, sw4 - 0.35, 2.0,
        font=SANS, size=11, color=MUTED)
    hline(s, cx + 0.2, cy + 3.45, sw4 - 0.35, color=RULE)
    txt(s, "OUTPUT", cx + 0.2, cy + 3.6, sw4 - 0.35, 0.22,
        font=MONO, size=7.5, color=MUTED)
    txt(s, output, cx + 0.2, cy + 3.86, sw4 - 0.35, 0.72,
        font=SERIF, size=13, color=INK, italic=True)

footer(s, "Longhorn × BridgeCom", "17 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 18 · SECTION: THE PROMISE
# ─────────────────────────────────────────────────────────────────────────────
divider_slide(prs, "05", "Promise",
              "What BridgeCom will commit to in writing, and what we hold ourselves to internally.")


# ─────────────────────────────────────────────────────────────────────────────
# 19 · KPI PHILOSOPHY
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Promise · 01")
eyebrow(s, "KPI philosophy", PAD_X, PAD_T)
txt(s, "We promise what the industry average delivers.\nWe plan to significantly exceed it.",
    PAD_X, PAD_T + 0.38, CW, 1.1,
    font=SERIF, size=40, color=INK, italic=True)

# Left body
txt(s, "Every KPI presented to Longhorn is set at or below global industry average benchmarks "
       "for financial services in emerging markets. This is deliberate. The targets we commit "
       "to in the contract are the floor we will not miss.\n\n"
       "An agency that over-promises and under-delivers loses the client at month three. "
       "An agency that under-promises and over-delivers earns trust, renewals, and referrals. "
       "The latter is the only commercially sustainable model.",
    PAD_X, PAD_T + 1.75, 5.8, 3.5,
    font=SANS, size=14, color=MUTED)

# Right boxes
boxes = [
    (WHITE,  RULE,  MUTED,  "THE FLOOR, IN THE CONTRACT",
     "What we put in writing and report against monthly. Calibrated to emerging-market averages. Credible. Conservative."),
    (NAVY,   None,  WHITE,  "THE CEILING, INSIDE THE TEAM",
     "What BridgeCom's team works to hit. Roughly 2x the floor on most metrics. Tracked weekly, never shared as a target."),
    (WHITE,  RED,   RED,    "IF WE MISS THE FLOOR",
     "After two consecutive misses, a written recovery plan presented before the third month begins. No hiding. A plan."),
]
bx = PAD_X + 6.2
for i, (bg_c, br_c, lbl_c, lbl, body) in enumerate(boxes):
    by = PAD_T + 1.65 + i * 1.5
    rect(s, bx, by, CW - 6.2, 1.35, fill=bg_c,
         line=br_c if br_c else None, line_w=Pt(1.2))
    tc = WHITE if bg_c == NAVY else MUTED
    txt(s, lbl, bx + 0.28, by + 0.16, CW - 6.6, 0.24,
        font=MONO, size=8, color=lbl_c)
    txt(s, body, bx + 0.28, by + 0.44, CW - 6.6, 0.76,
        font=SANS, size=12, color=tc)

footer(s, "Longhorn × BridgeCom", "19 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 20 · MONTH 12 COMMITMENTS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Promise · 02")
eyebrow(s, "What Month 12 looks like — committed numbers", PAD_X, PAD_T)
txt(s, "A year from signing, this is what we owe you.",
    PAD_X, PAD_T + 0.38, CW, 0.6,
    font=SERIF, size=38, color=INK, italic=True)

kpis = [
    (NAVY,  "TOTAL ENQUIRIES / MONTH",       "35",        " floor",  "Across all five channels. Stretch: 65. Source-attributed in CRM."),
    (RED,   "POLICIES / MONTH, DIGITAL",      "5–10",      "",        "Signed and bound. Stretch: 12–20. Attributed to channel and campaign."),
    (GOLD,  "MONTHLY PREMIUM REVENUE",        "K 80k",     "+",       "From digital-sourced clients. Stretch: K 200k+."),
    (GREEN, "WHATSAPP SUBSCRIBERS",           "700",       "+",       "Opt-in list. 40% open rate floor. Highest-converting audience you own."),
    (NAVY,  "SOCIAL FOLLOWERS, COMBINED",     "4.3k",      "+",       "FB 2.5k · IG 1.2k · LinkedIn 600. Stretch: ~8k. Real, Lusaka-led."),
    (RED,   "ACTIVE REFERRAL PARTNERS",       "4",         "",        "Signed agreements, monthly check-in. Referrals convert 2–3x cold leads."),
    (GOLD,  "KEYWORDS RANKING PAGE 1",        "3",         "+",       "Organic search on broker- and product-intent terms. SEO compounds."),
    (GREEN, "ENQUIRY RESPONSE SLA",           "<90s",      " auto",   "Automated WhatsApp reply. Human contact under 3h on A-leads. Day one."),
]
kw = CW / 4 - 0.1
kh = 2.38
for i, (color, label, val, suffix, ctx) in enumerate(kpis):
    col = i % 4
    row = i // 4
    kx = PAD_X + col * (kw + 0.13)
    ky = PAD_T + 1.2 + row * (kh + 0.12)
    rect(s, kx, ky, kw, kh, fill=WHITE, line=RULE, line_w=Pt(0.5))
    rect(s, kx, ky, kw, 0.05, fill=color)
    txt(s, label, kx + 0.18, ky + 0.18, kw - 0.3, 0.35,
        font=MONO, size=7.5, color=MUTED)
    txt(s, val + suffix, kx + 0.18, ky + 0.55, kw - 0.3, 0.85,
        font=SERIF, size=44, color=color)
    txt(s, ctx, kx + 0.18, ky + 1.42, kw - 0.3, 0.82,
        font=SANS, size=10, color=MUTED)

footer(s, "Longhorn × BridgeCom", "20 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 21 · SECTION: THE INVESTMENT
# ─────────────────────────────────────────────────────────────────────────────
divider_slide(prs, "06", "Investment",
              "Every cost. Transparent. Justified. Nothing hidden.")


# ─────────────────────────────────────────────────────────────────────────────
# 22 · YEAR-ONE BUDGET
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CREAM)
corner_lockup(s, "The Investment · 01")
eyebrow(s, "Year-one investment, all categories", PAD_X, PAD_T)
txt(s, "Every kwacha, accounted for.",
    PAD_X, PAD_T + 0.38, CW, 0.55,
    font=SERIF, size=42, color=INK, italic=True)

# Table header
ty = PAD_T + 1.15
hline(s, PAD_X, ty, CW, color=INK, width=Pt(1.2))
hdrs = [("CATEGORY", 4.5, PP_ALIGN.LEFT),
        ("ANNUAL LOW (ZMW)", 2.1, PP_ALIGN.RIGHT),
        ("ANNUAL HIGH (ZMW)", 2.1, PP_ALIGN.RIGHT),
        ("PAID TO", 3.0, PP_ALIGN.LEFT)]
hx = PAD_X
for hdr, hw, ha in hdrs:
    txt(s, hdr, hx + 0.05, ty + 0.06, hw - 0.08, 0.25,
        font=MONO, size=7.5, color=MUTED, align=ha)
    hx += hw + 0.06
hline(s, PAD_X, ty + 0.35, CW, color=INK, width=Pt(0.8))

budget_rows = [
    ("One-off setup (months 1–2)",            "K 50,500",    "K 78,200",    "BridgeCom Narratives Ltd"),
    ("Monthly retainer × 11 months (M2–M12)", "K 390,500",   "K 594,000",   "BridgeCom Narratives Ltd"),
    ("Technology & tools (annual)",           "K 15,624",    "K 31,260",    "Direct to platforms"),
    ("Advertising spend × 10 months (M3–M12)","K 117,600",   "K 243,600",   "Direct to LinkedIn, Meta, Google"),
    ("Video & production (as commissioned)",  "K 45,600",    "K 90,000",    "Local production partners"),
]
for r, (cat, lo, hi, paid) in enumerate(budget_rows):
    ry = ty + 0.35 + r * 0.68
    row_data = [(cat, 4.5, PP_ALIGN.LEFT, INK, False),
                (lo,  2.1, PP_ALIGN.RIGHT, INK, False),
                (hi,  2.1, PP_ALIGN.RIGHT, INK, False),
                (paid,3.0, PP_ALIGN.LEFT,  MUTED, False)]
    rx = PAD_X
    for text, rw, ra, rc, rb in row_data:
        txt(s, text, rx + 0.05, ry + 0.12, rw - 0.08, 0.44,
            font=SANS, size=12, color=rc, align=ra, bold=rb)
        rx += rw + 0.06
    hline(s, PAD_X, ry + 0.6, CW, color=RULE, width=Pt(0.4))

# Grand total row
gr_y = ty + 0.35 + 5 * 0.68
rect(s, PAD_X, gr_y, CW, 0.65, fill=NAVY)
totals = [("Year-one total range", 4.5, PP_ALIGN.LEFT),
          ("K 619,824",            2.1, PP_ALIGN.RIGHT),
          ("K 1,037,060",          2.1, PP_ALIGN.RIGHT),
          ("Full programme",       3.0, PP_ALIGN.LEFT)]
gx = PAD_X
for text, gw, ga in totals:
    txt(s, text, gx + 0.12, gr_y + 0.12, gw - 0.12, 0.42,
        font=SERIF, size=18, color=WHITE, align=ga)
    gx += gw + 0.06

txt(s, "Retainer minimum: six months. Monthly invoicing, 15-day terms. "
       "Ad spend invoiced separately and paid directly to platforms for full transparency. "
       "Longhorn owns every account, every contact, every asset, at all times.",
    PAD_X, gr_y + 0.82, CW, 0.42,
    font=SANS, size=11, color=MUTED)

footer(s, "Longhorn × BridgeCom", "22 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 23 · RETURN CASE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, CRMWRM)
corner_lockup(s, "The Investment · 02")
eyebrow(s, "The return, in plain numbers", PAD_X, PAD_T)
txt(s, "The programme is designed to pay for itself\nout of its own monthly output.",
    PAD_X, PAD_T + 0.38, CW, 1.0,
    font=SERIF, size=40, color=INK, italic=True)

roi_cards = [
    (WHITE, RULE,  INK,   "ALL-IN MONTHLY COST (AVG)",       "K 51k – 86k",
     "Year-one total divided by 12. Agency fees, tools, ad spend, production, everything."),
    (WHITE, RULE,  INK,   "MONTHLY PREMIUM, M12 (COMMITTED)", "K 80k+",
     "From digital-sourced clients at the contracted floor. Stretch target: K 200k+."),
    (NAVY,  None,  WHITE, "WHAT THAT MEANS",                  "Break-even",
     "By Month 12, the digital programme covers its own monthly cost. Everything beyond that is the return."),
]
rw = CW / 3 - 0.12
for i, (bg_c, br_c, tc, lbl, val, ctx) in enumerate(roi_cards):
    rx = PAD_X + i * (rw + 0.18)
    ry = PAD_T + 1.65
    rect(s, rx, ry, rw, 3.2, fill=bg_c, line=br_c if br_c else None, line_w=Pt(1))
    vc = WHITE if bg_c == NAVY else INK
    cc = RGBColor(0xC0,0xC0,0xE0) if bg_c == NAVY else MUTED
    txt(s, lbl, rx + 0.26, ry + 0.22, rw - 0.45, 0.28,
        font=MONO, size=8, color=RGBColor(0xA0,0xA0,0xD0) if bg_c == NAVY else MUTED)
    txt(s, val, rx + 0.26, ry + 0.6, rw - 0.45, 1.1,
        font=SERIF, size=48, color=vc)
    txt(s, ctx, rx + 0.26, ry + 1.85, rw - 0.45, 1.15,
        font=SANS, size=13, color=cc)

txt(s, "This is the conservative case. Stretch performance, which the team works to hit internally, "
       "delivers K 200,000+ in monthly premium revenue by M12 against the same cost base. "
       "Compounding into year two, the cost-per-enquiry falls 30%+ as organic and referral channels mature.",
    PAD_X, PAD_T + 5.1, CW, 0.65,
    font=SANS, size=13, color=MUTED)

footer(s, "Longhorn × BridgeCom", "23 / 24")


# ─────────────────────────────────────────────────────────────────────────────
# 24 · WHAT YOU OWN & NEXT STEPS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
corner_lockup(s, "Close", dark=True)
eyebrow(s, "Closing", PAD_X, PAD_T, dark=True)
txt(s, "Longhorn owns everything.\nBridgeCom holds nothing back.",
    PAD_X, PAD_T + 0.38, CW, 1.3,
    font=SERIF, size=58, color=WHITE, italic=True)

own_cards = [
    ("All accounts & admin rights",
     "Facebook, Instagram, LinkedIn, WhatsApp Business, Google Ads, Analytics, CRM. "
     "Longhorn is the primary owner. BridgeCom has agency access only."),
    ("All assets, data & contacts",
     "Website design files, every piece of content produced, the CRM contact database, "
     "the WhatsApp opt-in list. If the engagement ends, Longhorn keeps all of it."),
]
for i, (h, d) in enumerate(own_cards):
    ox = PAD_X + i * (CW / 2 + 0.12)
    oy = PAD_T + 2.0
    rect(s, ox, oy, CW / 2 - 0.12, 1.55,
         fill=RGBColor(0x18,0x10,0x90),
         line=RGBColor(0x50,0x50,0xB0), line_w=Pt(1))
    txt(s, h, ox + 0.26, oy + 0.18, CW / 2 - 0.55, 0.45,
        font=SERIF, size=22, color=WHITE)
    txt(s, d, ox + 0.26, oy + 0.68, CW / 2 - 0.55, 0.75,
        font=SANS, size=12, color=RGBColor(0xB0,0xB0,0xD8))

# Bottom section
hline(s, PAD_X, PAD_T + 4.0, CW, color=RGBColor(0x30,0x30,0x80), width=Pt(0.8))
txt(s, "NEXT STEP", PAD_X, PAD_T + 4.22, 3, 0.25,
    font=MONO, size=8, color=RGBColor(0x90,0x90,0xC0))
txt(s, "A two-hour working session with Longhorn leadership to walk every channel, "
       "every KPI and every line of the budget before signature.",
    PAD_X, PAD_T + 4.52, 7.5, 1.1,
    font=SERIF, size=30, color=WHITE)

txt(s, "BridgeCom Narratives", 9.5, PAD_T + 4.35, 3.5, 0.42,
    font=SANS, size=20, color=WHITE, bold=True)
txt(s, "hello@bridgecomnarratives.com", 9.5, PAD_T + 4.85, 3.5, 0.3,
    font=MONO, size=11, color=RGBColor(0x90,0x90,0xC0))
txt(s, "Lusaka, Zambia  ·  May 2026", 9.5, PAD_T + 5.2, 3.5, 0.3,
    font=MONO, size=11, color=RGBColor(0x90,0x90,0xC0))

footer(s, "Longhorn × BridgeCom", "24 / 24", dark=True)


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out = "/home/claude/repo/project/Longhorn-BridgeCom-Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
